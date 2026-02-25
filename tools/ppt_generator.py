#!/usr/bin/env python3
"""
PPT 生成工具 - 根据文章内容自动生成 PPT

用法:
  python tools/ppt_generator.py --title "标题" --content "内容" --output output.pptx
  python tools/ppt_generator.py --slides "幻灯片1内容|幻灯片2内容|..."
"""

import argparse
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# 配色方案
COLORS = {
    'primary': RGBColor(26, 54, 93),      # 深蓝色 #1A365D
    'secondary': RGBColor(45, 125, 210),   # 科技蓝 #2D7DD2
    'accent': RGBColor(237, 137, 54),     # 橙色 #ED8936
    'text': RGBColor(74, 85, 104),        # 深灰色 #4A5568
    'light_bg': RGBColor(247, 250, 252),  # 浅灰色 #F7FAFC
}


def create_title_slide(prs, title, subtitle=""):
    """创建标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 添加背景形状
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()

    # 添加标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 添加副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(1), Inches(4),
            Inches(8), Inches(1)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER


def create_content_slide(prs, title, content_lines):
    """创建内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 添加背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.2)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 内容
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(5.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['text']
        p.space_after = Pt(12)


def create_summary_slide(prs, title, points):
    """创建总结页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(1.2)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['secondary']
    bg.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 要点
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(5)
    )
    tf = content_box.text_frame

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS['text']
        p.space_after = Pt(16)


def create_closing_slide(prs, title="感谢聆听", subtitle=""):
    """创建结尾页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(3),
        Inches(8), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.5),
            Inches(8), Inches(1)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER


def generate_ppt(output_path, slides_data):
    """
    生成 PPT

    Args:
        output_path: 输出文件路径
        slides_data: 幻灯片数据列表
            [
                {"type": "title", "title": "标题", "subtitle": "副标题"},
                {"type": "content", "title": "标题", "content": ["内容1", "内容2"]},
                {"type": "summary", "title": "标题", "points": ["要点1", "要点2"]},
                {"type": "closing", "title": "标题", "subtitle": "副标题"}
            ]
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_data in slides_data:
        slide_type = slide_data.get("type", "content")

        if slide_type == "title":
            create_title_slide(
                prs,
                slide_data.get("title", ""),
                slide_data.get("subtitle", "")
            )
        elif slide_type == "content":
            create_content_slide(
                prs,
                slide_data.get("title", ""),
                slide_data.get("content", [])
            )
        elif slide_type == "summary":
            create_summary_slide(
                prs,
                slide_data.get("title", ""),
                slide_data.get("points", [])
            )
        elif slide_type == "closing":
            create_closing_slide(
                prs,
                slide_data.get("title", "感谢聆听"),
                slide_data.get("subtitle", "")
            )

    prs.save(output_path)
    print(f"[OK] PPT 已生成: {output_path}")


def main():
    parser = argparse.ArgumentParser(description="PPT 生成工具")
    parser.add_argument("--title", "-t", help="PPT 标题")
    parser.add_argument("--slides", "-s", help="幻灯片内容，用 | 分隔")
    parser.add_argument("--output", "-o", default="output.pptx", help="输出文件路径")
    parser.add_argument("--json", help="JSON 格式的幻灯片数据")

    args = parser.parse_args()

    if args.json:
        # 从 JSON 文件读取
        import json
        with open(args.json, 'r', encoding='utf-8') as f:
            slides_data = json.load(f)
        generate_ppt(args.output, slides_data)
        return

    if not args.slides:
        print("[ERROR] 请提供幻灯片内容")
        print("用法: python ppt_generator.py --title '标题' --slides '内容1|内容2' --output test.pptx")
        sys.exit(1)

    # 简单模式：标题页 + 内容页
    slides = [
        {"type": "title", "title": args.title or "演示文稿", "subtitle": ""},
    ]

    content_lines = args.slides.split("|")
    slides.append({
        "type": "content",
        "title": "内容",
        "content": content_lines
    })

    generate_ppt(args.output, slides)


if __name__ == "__main__":
    main()
